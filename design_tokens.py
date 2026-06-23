"""
Design tokens for Claude Design PowerPoint generation.
Single source of truth — import this in every generated script.

Usage:
    from design_tokens import *
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Slide dimensions (16:9 widescreen)
# ---------------------------------------------------------------------------
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

HEADER_H = Inches(1.2)   # Header bar height for Layout B/C/D
MARGIN_L = Inches(0.5)   # Left margin for body content
MARGIN_R = Inches(0.3)   # Right margin for body content
MARGIN_T = Inches(1.4)   # Top of body content (below header bar)

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
BG           = RGBColor(0xF8, 0xFA, 0xFC)  # Slide background (near white)
HEADER_BG    = RGBColor(0x1E, 0x29, 0x3B)  # Header bar / title slide bg
HEADER_TEXT  = RGBColor(0xFF, 0xFF, 0xFF)  # Text on header bg
BODY_TEXT    = RGBColor(0x33, 0x41, 0x55)  # Main body text
ACCENT       = RGBColor(0x25, 0x63, 0xEB)  # Accent / links / highlights
ACCENT_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)  # Light accent bg
CODE_BG      = RGBColor(0x0F, 0x17, 0x2A)  # Code block background
CODE_TEXT    = RGBColor(0xE2, 0xE8, 0xF0)  # Code text
DIVIDER      = RGBColor(0xCB, 0xD5, 0xE1)  # Horizontal / vertical dividers

# ---------------------------------------------------------------------------
# Typography
# ---------------------------------------------------------------------------
FONT_DEFAULT = "Arial"        # Slide title and body (cross-platform safe)
FONT_CODE    = "Courier New"  # Code blocks

SIZE_H1    = Pt(40)   # Slide title
SIZE_H2    = Pt(28)   # Section heading / subtitle
SIZE_BODY  = Pt(20)   # Body text / bullet points
SIZE_SMALL = Pt(16)   # Captions, footnotes
SIZE_CODE  = Pt(16)   # Code in code blocks

# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------

def set_slide_bg(slide, color=BG):
    """Set solid background color for a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, title_text, height=HEADER_H):
    """Add a dark header bar with title text (Layout B / C / D)."""
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = HEADER_BG
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.name = FONT_DEFAULT
    run.font.bold = True
    run.font.size = Pt(36)
    run.font.color.rgb = HEADER_TEXT


def add_text_box(slide, text, left, top, width, height,
                 font_name=FONT_DEFAULT, font_size=SIZE_BODY,
                 color=BODY_TEXT, bold=False, align=PP_ALIGN.LEFT):
    """Add a simple text box with consistent styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox


def add_code_block(slide, code_text, left, top, width, height):
    """Add a dark code block with monospace font."""
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.fill.background()

    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = code_text
    run.font.name = FONT_CODE
    run.font.size = SIZE_CODE
    run.font.color.rgb = CODE_TEXT


def new_blank_slide(prs):
    """Add a blank slide and set background color."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    return slide


def create_presentation():
    """Create a new 16:9 presentation."""
    from pptx import Presentation
    import os
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    os.makedirs("./output", exist_ok=True)
    return prs
